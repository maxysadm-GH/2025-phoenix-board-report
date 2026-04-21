"""Generate editable 3-slide PowerPoint deck from the Phoenix Board Report.
Uses proper tables for auto-sizing rows. All text is editable."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pathlib import Path

OUT = Path(__file__).parent

PURPLE = RGBColor(0x38, 0x16, 0x5F)
PLUM = RGBColor(0x3D, 0x11, 0x52)
LAVENDER = RGBColor(0xEA, 0xE1, 0xF5)
INK = RGBColor(0x1B, 0x1A, 0x18)
INK2 = RGBColor(0x4A, 0x47, 0x42)
MUTED = RGBColor(0x8A, 0x86, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
IVORY = RGBColor(0xFA, 0xF7, 0xF1)
HAIR = RGBColor(0xD9, 0xD3, 0xCB)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
PAD = Inches(0.5)


# ── helpers ──────────────────────────────────────────────

def _tx(slide, l, t, w, h, text, sz=11, color=INK, bold=False, italic=False, font='Calibri', align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    for attr, val in [('size', Pt(sz)), ('color.rgb', color), ('bold', bold), ('italic', italic), ('name', font)]:
        if '.' in attr:
            obj = r.font
            for part in attr.split('.')[:-1]: obj = getattr(obj, part)
            setattr(obj, attr.split('.')[-1], val)
        else:
            setattr(r.font, attr, val)
    return tf

def _rich(slide, l, t, w, h, parts, sz=11, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    for text, a in parts:
        r = p.add_run(); r.text = text
        r.font.size = Pt(a.get('sz', sz)); r.font.color.rgb = a.get('c', INK)
        r.font.bold = a.get('b', False); r.font.italic = a.get('i', False)
        r.font.name = a.get('f', 'Calibri')
    return tf

def _rect(slide, l, t, w, h, fill=None, border=None):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fill: s.fill.solid(); s.fill.fore_color.rgb = fill
    else: s.fill.background()
    if border: s.line.color.rgb = border; s.line.width = Pt(0.75)
    return s

def _no_border(tbl):
    for cell in tbl._tbl.iter_tcs():
        tcPr = cell.get_or_add_tcPr()
        for edge in ('lnL','lnR','lnT','lnB'):
            for old in tcPr.findall(qn(f'a:{edge}')): tcPr.remove(old)
            el = tcPr.makeelement(qn(f'a:{edge}'), {})
            el.append(tcPr.makeelement(qn('a:noFill'), {}))
            tcPr.append(el)

def _cell(c, text, sz=9, color=INK, bold=False, italic=False, font='Calibri', align=PP_ALIGN.LEFT):
    c.text = ''
    p = c.text_frame.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.color.rgb = color; r.font.bold = bold
    r.font.italic = italic; r.font.name = font
    c.margin_left = Pt(3); c.margin_right = Pt(3)
    c.margin_top = Pt(4); c.margin_bottom = Pt(4)

def _cell_rich(c, parts, sz=9, align=PP_ALIGN.LEFT):
    c.text = ''
    p = c.text_frame.paragraphs[0]; p.alignment = align
    for text, a in parts:
        r = p.add_run(); r.text = text
        r.font.size = Pt(a.get('sz', sz)); r.font.color.rgb = a.get('c', INK)
        r.font.bold = a.get('b', False); r.font.italic = a.get('i', False)
        r.font.name = a.get('f', 'Calibri')
    c.margin_left = Pt(3); c.margin_right = Pt(3)
    c.margin_top = Pt(4); c.margin_bottom = Pt(4)

def _stat(slide, l, t, num, label, nc=PURPLE, ns=24):
    _tx(slide, l, t, Inches(1.7), Inches(0.35), num, sz=ns, color=nc, font='Georgia')
    _tx(slide, l, t + Inches(0.33), Inches(1.7), Inches(0.3), label, sz=7, color=MUTED, bold=True)

def _bg(slide):
    slide.background.fill.solid(); slide.background.fill.fore_color.rgb = IVORY

def _masthead(slide, kicker, num):
    _tx(slide, PAD, Inches(0.22), Inches(9), Inches(0.2), kicker, sz=8, color=INK, bold=True)
    _rich(slide, Inches(10), Inches(0.22), Inches(2.8), Inches(0.2), [
        ('Fiscal 2025 Review ', {'sz': 8, 'c': MUTED}),
        (num, {'sz': 8, 'c': PURPLE, 'b': True}),
    ], align=PP_ALIGN.RIGHT)
    _rect(slide, PAD, Inches(0.42), Inches(12.3), Pt(0.75), fill=INK)


# ── SLIDE 1 ─────────────────────────────────────────────

def build_slide1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide)
    _masthead(slide, 'Vosges Haut-Chocolat \u00b7 A Conversation About Capability', 'I')

    _rich(slide, PAD, Inches(0.55), Inches(12.3), Inches(0.55), [
        ('Same walls. Same brand. ', {'sz': 28, 'f': 'Georgia'}),
        ('A different company ', {'sz': 28, 'f': 'Georgia', 'i': True}),
        ('underneath.', {'sz': 28, 'f': 'Georgia', 'i': True}),
    ])
    _tx(slide, PAD, Inches(1.05), Inches(12), Inches(0.2),
        'What the prior era built. What the new era built. And what we kept on purpose.',
        sz=10.5, color=INK2, italic=True, font='Georgia')

    # ── THEN ──
    LX, TY, LW = PAD, Inches(1.4), Inches(5.9)
    _rect(slide, LX, TY, LW, Inches(5.75), border=HAIR)
    _tx(slide, LX+Inches(0.2), TY+Inches(0.1), Inches(4), Inches(0.18),
        'Then \u00b7 Prior Ownership', sz=8, color=MUTED, bold=True)
    _rich(slide, LX+Inches(0.2), TY+Inches(0.3), Inches(5.4), Inches(0.25), [
        ('A brand ahead of its ', {'sz': 16, 'f': 'Georgia'}),
        ('engine.', {'sz': 16, 'f': 'Georgia', 'i': True}),
    ])
    _tx(slide, LX+Inches(0.2), TY+Inches(0.55), Inches(5.4), Inches(0.35),
        'World-class creative and innovation no one else could touch \u2014 but the operation was never built to carry the brand at scale.',
        sz=9, color=INK2, italic=True)

    rows = [
        ('Brand \u00b7 R&D', 'Iconic. Truffles, bars, bacon bar, PB-PB cup.'),
        ('Systems', '15-year-old CMS \u00b7 Nutracoster recipe tool, 20+ yrs \u00b7 islands of data.'),
        ('Floor tools', 'Paper pick lists. Paper counts. Rekeying everywhere.'),
        ('QA', 'Manual inspection. Errors caught downstream \u2014 often after shipment.'),
        ('Capacity', '~1,675 orders/day peak \u00b7 manual-heavy \u00b7 single carrier.'),
        ('Cost / ship', '$34.47 avg \u00b7 legacy pricing \u00b7 no carrier optionality.'),
        ('Service', 'Customer replies averaged 24\u201348 hours.'),
    ]
    ts = slide.shapes.add_table(len(rows), 2, LX+Inches(0.2), TY+Inches(0.95), Inches(5.5), Inches(2.1))
    t = ts.table; t.columns[0].width = Inches(0.95); t.columns[1].width = Inches(4.55); _no_border(t)
    for i,(k,v) in enumerate(rows):
        _cell(t.cell(i,0), k.upper(), sz=7.5, color=MUTED, bold=True)
        _cell(t.cell(i,1), v, sz=9, color=INK2)

    _tx(slide, LX+Inches(0.2), TY+Inches(3.2), Inches(5.4), Inches(0.35),
        '\u201cThe PB-PB cup never left a kitchen recipe.\u201d The creative was there. The chassis to ship it wasn\u2019t.',
        sz=9, color=INK2, italic=True, font='Georgia')

    sy = TY + Inches(3.7)
    _rect(slide, LX+Inches(0.2), sy, Inches(5.5), Pt(0.75), fill=INK)
    _stat(slide, LX+Inches(0.25), sy+Inches(0.08), '1,675', 'DAILY CAPACITY\nPEAK (2023)', nc=MUTED)
    _stat(slide, LX+Inches(2.1), sy+Inches(0.08), '$4.01', 'DIRECT FULFILLMENT\nLABOR COST', nc=MUTED)
    _stat(slide, LX+Inches(4.0), sy+Inches(0.08), '$34.47', 'COST PER\nSHIPMENT (AVG)', nc=MUTED)

    # ── NOW ──
    RX, RW = Inches(6.75), Inches(6.1)
    _rect(slide, RX, TY, RW, Inches(5.75), border=PURPLE)
    _tx(slide, RX+Inches(0.2), TY+Inches(0.1), Inches(4), Inches(0.18),
        'Now \u00b7 Current Ownership', sz=8, color=PURPLE, bold=True)
    _rich(slide, RX+Inches(0.2), TY+Inches(0.3), Inches(5.6), Inches(0.25), [
        ('The brand, ', {'sz': 16, 'f': 'Georgia', 'c': INK}),
        ('plus the engine it deserves.', {'sz': 16, 'f': 'Georgia', 'c': PURPLE, 'i': True}),
    ])
    _tx(slide, RX+Inches(0.2), TY+Inches(0.55), Inches(5.6), Inches(0.35),
        'Same building. Same team DNA. A modernized operational spine \u2014 connected systems, instrumented floors, simpler choices everywhere.',
        sz=9, color=INK2, italic=True)

    now_rows = [
        ('Brand \u00b7 R&D', [('Kept intact.', {'b':True,'c':PURPLE,'f':'Georgia','sz':9}),
                              (' Innovation calendar protected \u2014 the reason you buy Vosges.', {'sz':9})]),
        ('Systems',      [('Current-gen MRP', {'b':True,'c':PURPLE,'f':'Georgia','sz':9}),
                          (' retires the 15-yr CMS \u00b7 ', {'sz':9}),
                          ('Flavor Studio', {'b':True,'c':PURPLE,'f':'Georgia','sz':9}),
                          (' AI replaces Nutracoster.', {'sz':9})]),
        ('Floor tools',  [('Tablets & scanners', {'b':True,'c':PURPLE,'f':'Georgia','sz':9}),
                          (' on every station. Paper cut 50%.', {'sz':9})]),
        ('QA',           [('98.2% pick accuracy', {'b':True,'c':PURPLE,'f':'Georgia','sz':9}),
                          (' (Gorgias FY25) \u00b7 checkweigher + scan-verified.', {'sz':9})]),
        ('Capacity',     [('~3,120/day sustainable.', {'b':True,'c':PURPLE,'f':'Georgia','sz':9})]),
        ('Cost / ship',  [('$24.07 avg', {'b':True,'c':PURPLE,'f':'Georgia','sz':9}),
                          (' (FY24) | ', {'sz':9}),
                          ('$18.38 avg', {'b':True,'c':PURPLE,'f':'Georgia','sz':9}),
                          (' (FY25) \u2014 carrier simplification.', {'sz':9})]),
        ('Service',      [('Under 1-hour first-response', {'b':True,'c':PURPLE,'f':'Georgia','sz':9}),
                          (' with Gorgias + 24/7 virtual agents.', {'sz':9})]),
    ]
    ts2 = slide.shapes.add_table(len(now_rows), 2, RX+Inches(0.2), TY+Inches(0.95), Inches(5.7), Inches(2.1))
    t2 = ts2.table; t2.columns[0].width = Inches(0.95); t2.columns[1].width = Inches(4.75); _no_border(t2)
    for i,(k,parts) in enumerate(now_rows):
        _cell(t2.cell(i,0), k.upper(), sz=7.5, color=PLUM, bold=True)
        _cell_rich(t2.cell(i,1), parts, sz=9)

    ky = TY + Inches(3.2)
    _rect(slide, RX+Inches(0.2), ky, Inches(5.7), Inches(0.32), fill=LAVENDER)
    _rect(slide, RX+Inches(0.2), ky, Pt(3), Inches(0.32), fill=PURPLE)
    _rich(slide, RX+Inches(0.3), ky+Inches(0.04), Inches(5.5), Inches(0.25), [
        ('What we kept on purpose: ', {'b':True, 'c':PURPLE, 'sz':9}),
        ('the creative, the flavor R&D, the brand voice. None of the upgrade touched the product.', {'sz':9}),
    ])

    sy = TY + Inches(3.7)
    _rect(slide, RX+Inches(0.2), sy, Inches(5.7), Pt(0.75), fill=PURPLE)
    _stat(slide, RX+Inches(0.3), sy+Inches(0.08), '3,120', 'CAPACITY\n(SAME SPACE)')
    _stat(slide, RX+Inches(2.2), sy+Inches(0.08), '$1.38', 'DIRECT FULFILLMENT\nLABOR COST')
    _stat(slide, RX+Inches(4.2), sy+Inches(0.08), '$18.38', 'COST PER\nSHIPMENT (AVG)')


# ── SLIDE 2 ─────────────────────────────────────────────

def build_slide2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide)
    _masthead(slide, 'The Dividend of Modernization \u00b7 And What Held During the Rebuild', 'II')

    _rich(slide, PAD, Inches(0.55), Inches(12.3), Inches(0.55), [
        ('A footprint sized for December. ', {'sz': 28, 'f': 'Georgia'}),
        ('Running in April.', {'sz': 28, 'f': 'Georgia', 'i': True}),
    ])
    _tx(slide, PAD, Inches(1.05), Inches(12), Inches(0.2),
        'Fewer systems, better ones. More throughput, same walls. Faster answers at every door \u2014 even as the line was being rebuilt underneath.',
        sz=10.5, color=INK2, italic=True, font='Georgia')

    # 3 KPI tiles
    kpis = [('~3,120','Sustainable daily\nshipping capacity'),
            ('\u221250%','Technology cost\nyear over year'),
            ('\u221275%','Old software risk\neliminated')]
    kx = PAD
    for num, lbl in kpis:
        _rect(slide, kx, Inches(1.4), Inches(3.95), Pt(1.5), fill=INK)
        _tx(slide, kx+Inches(0.1), Inches(1.48), Inches(3.7), Inches(0.38), num,
            sz=24, color=PURPLE, font='Georgia')
        _tx(slide, kx+Inches(0.1), Inches(1.85), Inches(3.7), Inches(0.3), lbl,
            sz=7.5, color=MUTED, bold=True)
        kx += Inches(4.15)

    # ── LEFT: Modernization table ──
    _tx(slide, PAD, Inches(2.45), Inches(7), Inches(0.18),
        'What the Modernization Delivered', sz=8, color=PURPLE, bold=True)

    items = [
        ('MRP',         'Retired the 15-year CMS; current-gen MRP in its place.',            '1 stack'),
        ('Recipe',      'Flavor Studio replaces Nutracoster (20+ yrs) \u2014 AI at the bench.','Next-gen'),
        ('Floor',       'Tablets & scanners at every station. Paper \u221250%.',               'Every station'),
        ('Shipping',    'Carrier simplification. Ice & unboxing standardized.',               '\u221247% $/ship'),
        ('Service',     'Gorgias + 24/7 agents \u2014 answered under the hour.',              '24\u201348h \u2192 1h'),
        ('Visibility',  'Live dashboards. 24/7 automated audits. Data-driven decisions.',     'Always-on'),
        ('Integration', 'Every system talks to every other system.',                          'One graph'),
    ]
    ts = slide.shapes.add_table(len(items), 3, PAD, Inches(2.65), Inches(7.2), Inches(2.5))
    t = ts.table; t.columns[0].width=Inches(0.9); t.columns[1].width=Inches(4.9); t.columns[2].width=Inches(1.4)
    _no_border(t)
    for i,(k,v,s) in enumerate(items):
        _cell(t.cell(i,0), k.upper(), sz=7.5, color=MUTED, bold=True)
        _cell(t.cell(i,1), v, sz=9)
        _cell(t.cell(i,2), s, sz=9, color=PURPLE, italic=True, font='Georgia', align=PP_ALIGN.RIGHT)

    # ── RIGHT: Pick accuracy ──
    rc = Inches(8.0)
    _rect(slide, rc, Inches(2.45), Inches(4.85), Inches(4.65), border=HAIR)

    _tx(slide, rc+Inches(0.25), Inches(2.55), Inches(4.3), Inches(0.18),
        'Pick Accuracy \u00b7 Gorgias FY25', sz=8, color=MUTED, bold=True)

    _rich(slide, rc+Inches(0.25), Inches(2.85), Inches(3), Inches(0.8), [
        ('98.2', {'sz':52, 'c':PURPLE, 'f':'Georgia'}),
        ('%', {'sz':26, 'c':MUTED, 'f':'Georgia'}),
    ])

    _tx(slide, rc+Inches(0.25), Inches(3.65), Inches(4.3), Inches(0.2),
        'of shipped orders \u2014 no customer-reported issue', sz=9, color=INK2, italic=True, font='Georgia')

    _rect(slide, rc+Inches(0.25), Inches(4.1), Pt(2), Inches(0.7), fill=PURPLE)
    _rich(slide, rc+Inches(0.35), Inches(4.12), Inches(4.2), Inches(0.65), [
        ('Held the line during the rebuild. ', {'b':True, 'sz':9.5}),
        ('Q4 2025 \u2014 the quarter we re-sequenced the pick line \u2014 shipped record holiday volume (Dec 2025: 61,503 shipments) without degrading the accuracy floor.',
         {'sz':9.5, 'i':True, 'f':'Georgia'}),
    ])

    # Capacity chart description (since Chart.js can't render in PPTX)
    _tx(slide, PAD, Inches(5.35), Inches(7), Inches(0.18),
        'Exhibit \u00b7 Capacity Already Built', sz=8, color=PURPLE, bold=True)
    _rich(slide, PAD, Inches(5.55), Inches(12), Inches(0.2), [
        ('Engineered to sustain ', {'sz':11, 'f':'Georgia', 'c':INK}),
        ('~3,120/day', {'sz':11, 'f':'Georgia', 'c':PURPLE, 'i':True}),
        ('. Same building. Same racking. Same crew. December 2025 averaged 2,563/day \u2014 on the same footprint that capped at 1,675/day in 2023.',
         {'sz':11, 'f':'Georgia', 'c':INK2}),
    ])

    # Monthly throughput table (replaces the chart)
    months = ['Apr','May','Jun','Jul','Aug','Sep','Oct','Nov','Dec','Jan','Feb','Mar']
    vals =   ['205','222','152','78','166','217','204','1,472','2,563','319','457','191']
    ts3 = slide.shapes.add_table(2, 12, PAD, Inches(5.9), Inches(12.3), Inches(0.55))
    t3 = ts3.table
    for i in range(12): t3.columns[i].width = Inches(1.025)
    _no_border(t3)
    for i in range(12):
        _cell(t3.cell(0,i), months[i], sz=7.5, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
        _cell(t3.cell(1,i), vals[i]+'/d', sz=9, color=PURPLE, font='Georgia', align=PP_ALIGN.CENTER)

    _tx(slide, PAD, Inches(6.5), Inches(12), Inches(0.18),
        'FY25 monthly daily average throughput (ShipStation) \u00b7 Sustainable ceiling: ~3,120/day',
        sz=7.5, color=MUTED)


# ── SLIDE 3 ─────────────────────────────────────────────

def build_slide3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); _bg(slide)
    _masthead(slide, 'What Lies Ahead \u00b7 Phoenix Phase Two \u00b7 And Beyond', 'III')

    _rich(slide, PAD, Inches(0.55), Inches(12.3), Inches(0.55), [
        ('Faster. Simpler. ', {'sz': 28, 'f': 'Georgia'}),
        ('And then \u2014 portable.', {'sz': 28, 'f': 'Georgia', 'i': True}),
    ])
    _tx(slide, PAD, Inches(1.05), Inches(12), Inches(0.2),
        'An $80,000 investment that pays itself back in twelve weeks \u2014 and teaches us how to wire the next warehouse in a week, not a year.',
        sz=10.5, color=INK2, italic=True, font='Georgia')

    # 6 KPI tiles
    kpis = [('$80K','Phase 2\ninvest'), ('$7K/wk','Ongoing\nsavings'),
            ('12 wks','Payback\nperiod'), ('$364K','Year-1\nannualized'),
            ('$1.1M','Three-year\ncumulative'), ('25\u201335%','Overhead cut\non the line')]
    kx = PAD
    for num, lbl in kpis:
        _rect(slide, kx, Inches(1.4), Inches(1.95), Pt(1.5), fill=INK)
        _tx(slide, kx+Inches(0.08), Inches(1.47), Inches(1.8), Inches(0.33), num,
            sz=18, color=PURPLE, font='Georgia')
        _tx(slide, kx+Inches(0.08), Inches(1.78), Inches(1.8), Inches(0.25), lbl,
            sz=7, color=MUTED, bold=True)
        kx += Inches(2.07)

    # ── LEFT: Phase Two ──
    _tx(slide, PAD, Inches(2.3), Inches(4.5), Inches(0.18),
        'Phase Two \u00b7 On the Line', sz=8, color=PURPLE, bold=True)
    _rich(slide, PAD, Inches(2.5), Inches(4.5), Inches(0.25), [
        ('The things that still look ', {'sz':14, 'f':'Georgia'}),
        ('manual.', {'sz':14, 'f':'Georgia', 'c':PURPLE, 'i':True}),
    ])

    phase2 = [
        ('Gift orders','Dedicated box-maker station at the line\u2019s head. Gift-note printed on scan \u2014 never missed.'),
        ('QA stop','QA & Organizer merge into one role. One person owns the box before it moves.'),
        ('Sealing','Auto-sealer replaces manual tape & fold. Seasonal rental \u2014 paid back in one peak week.'),
        ('Labeling','ShipStation Mobile retired. Auto peel-and-apply labeler takes over.'),
        ('Boxes','Sizes re-aligned to FedEx One Rate tiers. New SKUs require box-team approval.'),
        ('Bedding','Pre-cut, known-weight. Shipping math becomes exact.'),
    ]
    ts = slide.shapes.add_table(len(phase2), 2, PAD, Inches(2.8), Inches(4.5), Inches(2.5))
    t = ts.table; t.columns[0].width=Inches(0.85); t.columns[1].width=Inches(3.65); _no_border(t)
    for i,(tag,desc) in enumerate(phase2):
        _cell(t.cell(i,0), tag, sz=9, color=PURPLE, italic=True, font='Georgia')
        _cell(t.cell(i,1), desc, sz=9)

    # ── MIDDLE: The math ──
    mx = Inches(5.3)
    _tx(slide, mx, Inches(2.3), Inches(3.8), Inches(0.18),
        'The Math \u00b7 Cumulative $ Through the Door', sz=8, color=PURPLE, bold=True)
    _rich(slide, mx, Inches(2.5), Inches(3.8), Inches(0.25), [
        ('$80K out. ', {'sz':14, 'f':'Georgia'}),
        ('$1.1M back.', {'sz':14, 'f':'Georgia', 'c':PURPLE, 'i':True}),
    ])

    mh = [('Wk 12','Break-even'), ('Wk 52','+$284K net'), ('Yr 3','+$1M net')]
    mhx = mx
    for num, lbl in mh:
        _rect(slide, mhx, Inches(2.82), Inches(1.15), Pt(0.75), fill=INK)
        _tx(slide, mhx, Inches(2.88), Inches(1.15), Inches(0.28), num, sz=16, color=PURPLE, font='Georgia')
        _tx(slide, mhx, Inches(3.14), Inches(1.15), Inches(0.18), lbl.upper(), sz=7, color=MUTED, bold=True)
        mhx += Inches(1.25)

    # Payback table (replaces chart)
    weeks = ['W0','W4','W8','W12','W16','W26','W52','Yr 3']
    cumul = ['\u2212$80K','\u2212$52K','\u2212$24K','+$4K','+$32K','+$102K','+$284K','+$1.0M']
    ts2 = slide.shapes.add_table(2, len(weeks), mx, Inches(3.45), Inches(3.8), Inches(0.5))
    t2 = ts2.table
    for i in range(len(weeks)): t2.columns[i].width = Inches(3.8/len(weeks))
    _no_border(t2)
    for i in range(len(weeks)):
        _cell(t2.cell(0,i), weeks[i], sz=7, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
        c = PURPLE if not cumul[i].startswith('\u2212') else MUTED
        _cell(t2.cell(1,i), cumul[i], sz=8, color=c, font='Georgia', align=PP_ALIGN.CENTER)

    _rect(slide, mx, Inches(4.1), Pt(2), Inches(0.65), fill=PURPLE)
    _rich(slide, mx+Inches(0.1), Inches(4.12), Inches(3.6), Inches(0.6), [
        ('Read it this way \u2014 ', {'b':True, 'sz':9}),
        ('every line of cost stays flat; every line of savings stacks. The curves cross at week twelve. Everything after is structural margin on the P&L.',
         {'sz':9, 'i':True, 'f':'Georgia'}),
    ])

    # ── RIGHT: Vision ──
    vx = Inches(9.5)
    _rect(slide, vx, Inches(2.3), Inches(3.55), Inches(4.8), fill=PLUM)
    _tx(slide, vx+Inches(0.2), Inches(2.4), Inches(3.1), Inches(0.18),
        'And Then \u2014 The Real Unlock', sz=8, color=LAVENDER, bold=True)
    _rich(slide, vx+Inches(0.2), Inches(2.6), Inches(3.1), Inches(0.25), [
        ('A template that ', {'sz':14, 'f':'Georgia', 'c':WHITE}),
        ('travels.', {'sz':14, 'f':'Georgia', 'c':LAVENDER, 'i':True}),
    ])

    _tx(slide, vx+Inches(0.2), Inches(2.95), Inches(3.1), Inches(0.2),
        'Chicago (Hub)  \u2500\u2500  Node II  \u2500\u2500  Node III', sz=9, color=WHITE, font='Georgia')

    for txt, y in [
        ('Every simplification on this line is a step toward a standard \u2014 something we can wire somewhere else.', Inches(3.3)),
        ('A portable warehouse template. Light up a second location \u2014 closer to the customer, faster to deliver, cheaper to ship.', Inches(3.85)),
        ('Not a bigger building. A smarter network.', Inches(4.5)),
    ]:
        _tx(slide, vx+Inches(0.2), y, Inches(3.1), Inches(0.5), txt, sz=9.5, color=WHITE, font='Georgia')

    _rect(slide, vx+Inches(0.2), Inches(4.9), Inches(3.1), Pt(0.75), fill=LAVENDER)
    _rich(slide, vx+Inches(0.2), Inches(5.0), Inches(3.1), Inches(0.45), [
        ('Simplify first. Replicate second. ', {'f':'Georgia','sz':9.5,'c':LAVENDER,'i':True,'b':True}),
        ('The next warehouse is a configuration, not a construction project.', {'f':'Georgia','sz':9.5,'c':WHITE,'i':True}),
    ])


def build(name):
    prs = Presentation()
    prs.slide_width = SLIDE_W; prs.slide_height = SLIDE_H
    build_slide1(prs); build_slide2(prs); build_slide3(prs)
    out = OUT / f"{name}.pptx"
    prs.save(str(out))
    print(f"[{name}] -> {out}")

if __name__ == '__main__':
    build('editorial'); build('brand')
    print("Done.")
